import tkinter as tk
from pptx import Presentation
from pptx.util import Cm, Pt

# data = open(r"E:\测试\txt.txt")
# lines = len(data.readlines())
# data.read()
# print(data.read())
# print(lines)
# data.close()

with open(r"E:\测试\txt.txt", "r") as f:
  data = f.read()
  # print(data)

# 加载库
import  os
from pptx import Presentation

# 设置路径
work_path = r'E:\测试'
os.chdir(work_path)

# 实例化 ppt 文档对象
prs = Presentation()

# 选择布局
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
# 容器
shapes = slide.shapes
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame

## 分三行写入正文
### 正文文本
# d = len(data)//3 + 1
# ts = list(data)
# for i in range(10):
#     ts.insert((i+1)*d - 2 + i, '\n')  # 插入换行符
# s = ''.join(ts)
### 写入正文
new_para = tf.add_paragraph()  # 添加段落
new_para.text = data
# new_para.line_spacing = 1.5    # 1.5 倍的行距


# 保存 ppt
prs.save('test.pptx')

# ## 分三行写入正文
# ### 正文文本
# s ='你侬我侬，忒煞情多；情多处，热如火；把一块泥，捻一个你，塑一个我，将咱两个一齐打碎，用水调和；\
# 再捻一个你，再塑一个我。我泥中有你，你泥中有我；我与你生同一个衾，死同一个椁。'
# d = len(s)//3 + 1
# ts = list(s)
# for i in range(3):
#     ts.insert((i+1)*d - 2 + i, '\n')  # 插入换行符
# s = ''.join(ts)
# ### 写入正文
# new_para = tf.add_paragraph()  # 添加段落
# new_para.text = s
# new_para.line_spacing = 1.5    # 1.5 倍的行距
#
#
# # 保存 ppt
# prs.save('test.pptx')

