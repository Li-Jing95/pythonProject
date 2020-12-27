
from pptx import Presentation
from datetime import datetime

prs = Presentation("E:\测试\示例文件2.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])  # 用第一个母版生成一页ppt
for shape in slide.placeholders:         # 获取这一页所有的占位符
    phf = shape.placeholder_format
    print(f'{phf.idx}--{shape.name}--{phf.type}')
    shape.text = f'{phf.idx}--{phf.type}'  #在占位符中填写“占位符id号--占位符类型”
prs.save('E:\测试\示例文件2.pptx')