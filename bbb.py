import os
import tkinter as tk
from pptx import Presentation
from pptx.util import Cm, Pt


root = tk.Tk()
root.geometry('400x200')
root.title('选择检测规则')
# 第一行（两列）
row1 = tk.Frame(root)
row1.pack(fill="x")
l1 = tk.Label(row1, text='关键字信息：', width=12, height=2).pack(side=tk.LEFT)
root.name = tk.StringVar()
u1 = tk.Entry(row1, textvariable=root.name, width=20)
u1.pack()
# 第二行
row2 = tk.Frame(root)
row2.pack(fill="x")
l2 = tk.Label(row2, text='阈值：', width=12, height=2).pack(side=tk.LEFT)
root.name = tk.StringVar()
u2 = tk.Entry(row2, textvariable=root.name, width=20)
u2.pack()


def get():
    # 关键字信息
    data1 = u1.get()
    # 阈值
    data2 = u2.get()
    print(data1, data2)

    count = 1
    base_path = "E:/测试/"
    txt = open(base_path + "txt.txt","w")
    xls = open(base_path + "1.xls", "w")
    doc = open(base_path+"1.doc", "w")
    xlsx = open(base_path + "1.xlsx", "w")
    docx = open(base_path+"1.docx", "w",encoding='GB2312')

    while count <= int(data2):
        txt.write(data1 + '\n')
        xls.write(data1 + '\n')
        doc.write(data1 + '\n')
        xlsx.write(data1 + '\n')
        docx.write(data1 + '\n')
        count += 1
    xls.close()
    with open(r"E:\测试\txt.txt", "r") as f:
        data = f.read()

    # 加载库
    import os
    from pptx import Presentation

    # 设置路径
    work_path = r'E:\测试'
    os.chdir(work_path)

    # 实例化 ppt 文档对象
    prs = Presentation()

    # 选择布局
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    # 容器
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    ### 写入正文
    new_para = tf.add_paragraph()  # 添加段落
    new_para.text = data

    # 保存 ppt和pptx
    prs.save('1.pptx')
    prs.save('1.ppt')

tk.Button(root, text='确定', command=get).pack()
root.mainloop()
