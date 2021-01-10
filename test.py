import os
import os.path
import time
import tkinter as tk
import xlsxwriter
import xlwt
from pptx import Presentation
from docx import Document
from win32com.client import constants, gencache

root = tk.Tk()
root.geometry('520x200')
root.resizable(width=False, height=False)
root.title('选择检测规则')

# 第一行（两列）
row1 = tk.Frame(root)
row1.pack(fill="x")
l1 = tk.Label(row1, text='关键字信息：', width=12, height=2).pack(side=tk.LEFT)
root.name = tk.StringVar()
u1 = tk.Entry(row1, textvariable=root.name, width=20)
u1.pack()
# 第二行
row2 = tk.Frame(root)
row2.pack(fill="x")
l2 = tk.Label(row2, text='阈值：', width=12, height=2).pack(side=tk.LEFT)
root.name = tk.StringVar()
u2 = tk.Entry(row2, textvariable=root.name, width=20)
u2.pack()
# 第三行
row2 = tk.Frame(root)
row2.pack(fill="x")
l2 = tk.Label(row2, text='1、文件默认保存路径为 C:\\test，请在C盘下创建此文件夹，并在此文件夹下创建pptx.pptx文件', fg='red').pack(side=tk.LEFT)
# 第四行
row3 = tk.Frame(root)
row3.pack(fill="x")
l3 = tk.Label(row3, text='2、点击“确定”后直接关闭此窗口(此窗口默认5秒关闭)', fg='red').pack(side=tk.LEFT)

lbTime = tk.Label(root, fg='red', anchor='w')
lbTime.place(x=10, y=250, width=150)

def get():
    # 关键字信息
    data1 = u1.get()
    # 阈值
    data2 = u2.get()
    print(data1, data2)

    count = 1
    base_path = r"C:/test/"
    txt = open(base_path + "1.txt", "w")
    # xls = open(base_path + "1.xls", "w")
    # doc = open(base_path + "1.doc", "w")

    while count <= int(data2):
        txt.write(data1 + '\n')
        # xls.write(data1 + '\n')
        # doc.write(data1 + '\n')
        count += 1
    txt.close()

    with open(r"C:\test\1.txt", "r") as f:
        data = f.read()
    # ####################################
    # 写入ppt/pptx
    # 设置路径
    work_path = r'C:\test'
    os.chdir(work_path)
    #
    # 实例化ppt文档对象
    prs = Presentation('pptx.pptx')

    # 选择布局
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)

    # 容器
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    # 写入正文
    new_para = tf.add_paragraph()  # 添加段落
    new_para.text = data

    # 保存ppt和pptx
    prs.save('1.ppt')
    prs.save('1.pptx')
    prs.save('1.wpp')
    # #####################################
    #
    # 写入doc/docx文件
    doc = Document()
    p = doc.add_paragraph(data)
    doc.save('1.doc')
    doc.save('1.docx')
    doc.save('1.wps')
    ######################################

    # 写入xls/xlsx文件
    txtopen = open("C:/test/1.txt", 'r')
    lines = txtopen.readlines()
    # 新建一个excel文件
    xls = xlwt.Workbook(encoding='utf-8', style_compression=0)
    xlsx = xlsxwriter.Workbook(base_path + '1.xlsx')

    # 新建一个sheet
    xlssheet = xls.add_sheet('sheet1')
    xlsxsheet = xlsx.add_worksheet('sheet1')
    # 写入a.txt，a.txt文件有N行文件
    i = 0
    for line in lines:
        xlssheet.write(i, 0, line)
        xlsxsheet.write(i, 0, line)
        i = i + 1
    xls.save(base_path + '1.xls')
    xls.save(base_path + '1.et')
    xlsx.close()

    #####################################
    # 写入pdf文件
    docx_file_path = base_path + '1.docx'
    pdf_file_path = base_path + '1.pdf'
    w = gencache.EnsureDispatch('Word.Application')
    doc = w.Documents.Open(docx_file_path, ReadOnly=1)
    doc.ExportAsFixedFormat(pdf_file_path,
                            constants.wdExportFormatPDF,
                            Item=constants.wdExportDocumentWithMarkup,
                            CreateBookmarks=constants.wdExportCreateHeadingBookmarks)

    w.Quit(constants.wdDoNotSaveChanges)
    time.sleep(5)
    root.destroy()
    # def autoClose():
    #     for i in range(5):
    #         lbTime['text'] = '距离窗口关闭还有{}秒'.format(5 - i)
    #         time.sleep(1)
    #     root.destroy()
    # t = threading.Thread(target=autoClose())
    # t.start()

tk.Button(root, text='确定', command=get).pack()
root.mainloop()